from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

import pytest
from pptx import Presentation
from pptx.util import Inches

from quarry.models import PageType
from quarry.presentation_processor import (
    SUPPORTED_PRESENTATION_EXTENSIONS,
    _extract_slide_text,
    _format_slide_content,
    _table_to_latex,
    process_presentation_file,
)

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType
    from pptx.slide import Slide
    from pptx.table import Table


def _make_pptx(tmp_path: Path, name: str = "test.pptx") -> Path:
    """Return the path for a new PPTX file."""
    return tmp_path / name


def _new_prs() -> PresentationType:
    return Presentation()


def _save(prs: PresentationType, path: Path) -> None:
    prs.save(str(path))


def _add_table(
    slide: Slide,
    rows: int,
    cols: int,
    *,
    left: float = 1,
    top: float = 1,
    width: float = 4,
    height: float = 2,
) -> Table:
    """Add a table to a slide, returning the Table object."""
    shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    return shape.table


class TestSupportedExtensions:
    def test_includes_pptx(self):
        assert ".pptx" in SUPPORTED_PRESENTATION_EXTENSIONS

    def test_does_not_include_ppt(self):
        assert ".ppt" not in SUPPORTED_PRESENTATION_EXTENSIONS

    def test_no_overlap_with_other_extensions(self):
        from quarry.code_processor import SUPPORTED_CODE_EXTENSIONS
        from quarry.html_processor import SUPPORTED_HTML_EXTENSIONS
        from quarry.spreadsheet_processor import (
            SUPPORTED_SPREADSHEET_EXTENSIONS,
        )
        from quarry.text_processor import SUPPORTED_TEXT_EXTENSIONS

        overlap = SUPPORTED_PRESENTATION_EXTENSIONS & (
            SUPPORTED_CODE_EXTENSIONS
            | SUPPORTED_TEXT_EXTENSIONS
            | SUPPORTED_SPREADSHEET_EXTENSIONS
            | SUPPORTED_HTML_EXTENSIONS
        )
        assert overlap == frozenset(), f"Overlapping extensions: {overlap}"


class TestTableToLatex:
    def test_basic_table(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = _add_table(slide, 3, 2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Age"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "30"
        table.cell(2, 0).text = "Bob"
        table.cell(2, 1).text = "25"

        result = _table_to_latex(table)

        assert r"\begin{tabular}" in result
        assert "Name & Age" in result
        assert "Alice & 30" in result
        assert "Bob & 25" in result

    def test_empty_table(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = _add_table(slide, 1, 1, width=2, height=1)
        table.cell(0, 0).text = ""

        result = _table_to_latex(table)

        assert r"\begin{tabular}" in result

    def test_special_chars_escaped(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = _add_table(slide, 2, 1, width=2, height=1)
        table.cell(0, 0).text = "Price"
        table.cell(1, 0).text = "$100"

        result = _table_to_latex(table)

        assert r"\$100" in result


class TestExtractSlideText:
    def test_title_and_body(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "My Title"
        slide.placeholders[1].text = "Subtitle text"

        title, body, _notes = _extract_slide_text(slide)

        assert title == "My Title"
        assert "Subtitle text" in body

    def test_speaker_notes(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txbox.text_frame.text = "Body content"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "These are speaker notes."

        _title, _body, notes = _extract_slide_text(slide)

        assert notes == "These are speaker notes."

    def test_no_notes(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        _title, _body, notes = _extract_slide_text(slide)

        assert notes == ""

    def test_table_in_slide(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table = _add_table(slide, 2, 2)
        table.cell(0, 0).text = "X"
        table.cell(0, 1).text = "Y"
        table.cell(1, 0).text = "1"
        table.cell(1, 1).text = "2"

        _title, body, _notes = _extract_slide_text(slide)

        assert r"\begin{tabular}" in body
        assert "X & Y" in body

    def test_empty_slide(self):
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title, body, notes = _extract_slide_text(slide)

        assert title == ""
        assert body == ""
        assert notes == ""


class TestFormatSlideContent:
    def test_full_content(self):
        result = _format_slide_content("Title", "Body text", "Notes here")

        assert result.startswith("# Title")
        assert "Body text" in result
        assert "---\nSpeaker Notes:" in result
        assert "Notes here" in result

    def test_no_title(self):
        result = _format_slide_content("", "Body only", "")

        assert not result.startswith("#")
        assert "Body only" in result

    def test_no_notes(self):
        result = _format_slide_content("Title", "Body", "")

        assert "Speaker Notes" not in result

    def test_no_body(self):
        result = _format_slide_content("Title", "", "Notes")

        assert "# Title" in result
        assert "Speaker Notes" in result
        assert "Notes" in result

    def test_all_empty(self):
        result = _format_slide_content("", "", "")
        assert result == ""


class TestProcessPresentationFile:
    def test_basic_pptx(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide One"
        slide.placeholders[1].text = "Content here"
        _save(prs, f)

        pages = process_presentation_file(f)

        assert len(pages) == 1
        assert pages[0].page_type == PageType.PRESENTATION
        assert pages[0].document_name == "test.pptx"
        assert "Slide One" in pages[0].text
        assert "Content here" in pages[0].text

    def test_multiple_slides(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f"Slide {i + 1}"
            slide.placeholders[1].text = f"Content {i + 1}"
        _save(prs, f)

        pages = process_presentation_file(f)

        assert len(pages) == 3
        for i, page in enumerate(pages):
            assert page.page_number == i + 1
            assert page.total_pages == 3
            assert f"Slide {i + 1}" in page.text

    def test_empty_slides_skipped(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Real Slide"
        slide1.placeholders[1].text = "Has content"
        prs.slides.add_slide(prs.slide_layouts[5])  # blank
        slide3 = prs.slides.add_slide(prs.slide_layouts[0])
        slide3.shapes.title.text = "Another Slide"
        slide3.placeholders[1].text = "More content"
        _save(prs, f)

        pages = process_presentation_file(f)

        assert len(pages) == 2
        assert pages[0].page_number == 1
        assert pages[1].page_number == 2
        assert pages[0].total_pages == 2
        assert "Real Slide" in pages[0].text
        assert "Another Slide" in pages[1].text

    def test_slide_with_table(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(4), Inches(0.5))
        txbox.text_frame.text = "Revenue Data"
        table = _add_table(slide, 2, 2, top=1.5)
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Sales"
        table.cell(1, 0).text = "North"
        table.cell(1, 1).text = "1000"
        _save(prs, f)

        pages = process_presentation_file(f)

        assert len(pages) == 1
        assert "Revenue Data" in pages[0].text
        assert r"\begin{tabular}" in pages[0].text
        assert "Region & Sales" in pages[0].text
        assert "North & 1000" in pages[0].text

    def test_slide_with_notes(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Body"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Remember to mention X"
        _save(prs, f)

        pages = process_presentation_file(f)

        assert len(pages) == 1
        assert "---\nSpeaker Notes:" in pages[0].text
        assert "Remember to mention X" in pages[0].text

    def test_document_name_default(self, tmp_path: Path):
        f = _make_pptx(tmp_path, "deck.pptx")
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title"
        slide.placeholders[1].text = "Body"
        _save(prs, f)

        pages = process_presentation_file(f)

        assert pages[0].document_name == "deck.pptx"

    def test_document_name_override(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title"
        slide.placeholders[1].text = "Body"
        _save(prs, f)

        pages = process_presentation_file(f, document_name="subdir/deck.pptx")

        assert pages[0].document_name == "subdir/deck.pptx"

    def test_document_path_is_resolved(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title"
        slide.placeholders[1].text = "Body"
        _save(prs, f)

        pages = process_presentation_file(f)

        assert pages[0].document_path == str(f.resolve())

    def test_all_empty_slides(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        prs.slides.add_slide(prs.slide_layouts[5])
        prs.slides.add_slide(prs.slide_layouts[5])
        _save(prs, f)

        pages = process_presentation_file(f)

        assert pages == []

    def test_special_chars_in_title_escaped(self, tmp_path: Path):
        f = _make_pptx(tmp_path)
        prs = _new_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Revenue: $4.2M & Growth"
        slide.placeholders[1].text = "Details"
        _save(prs, f)

        pages = process_presentation_file(f)

        assert r"\$4.2M" in pages[0].text
        assert r"\&" in pages[0].text


class TestProcessPresentationErrors:
    def test_unsupported_extension(self, tmp_path: Path):
        f = tmp_path / "data.ppt"
        f.write_bytes(b"\x00")

        with pytest.raises(ValueError, match="Unsupported presentation format"):
            process_presentation_file(f)

    def test_unsupported_extension_odp(self, tmp_path: Path):
        f = tmp_path / "data.odp"
        f.write_bytes(b"\x00")

        with pytest.raises(ValueError, match="Unsupported presentation format"):
            process_presentation_file(f)
